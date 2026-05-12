"""
Generate MODELS.pptx from MODELS.md.

Equations are rendered to PNG via matplotlib + LaTeX; tables are native
PowerPoint tables; paragraphs are editable text boxes.
"""
import hashlib
import os
import shutil
import tempfile

os.environ.setdefault("MPLCONFIGDIR", "/tmp/claude/mpl")

import matplotlib

matplotlib.use("Agg")
matplotlib.rcParams["text.usetex"] = True
matplotlib.rcParams["text.latex.preamble"] = (
    r"\usepackage{amsmath}\usepackage{amssymb}"
)
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PPTX = os.path.join(HERE, "MODELS.pptx")
EQ_CACHE = tempfile.mkdtemp(prefix="eqs_")

# Palette: Ocean Gradient
COL_BG = RGBColor(0xFF, 0xFF, 0xFF)
COL_TITLE = RGBColor(0x06, 0x5A, 0x82)     # deep blue
COL_ACCENT = RGBColor(0x1C, 0x72, 0x93)    # teal
COL_DARK = RGBColor(0x21, 0x29, 0x5C)      # midnight
COL_BODY = RGBColor(0x1A, 0x1A, 0x1A)
COL_MUTED = RGBColor(0x55, 0x55, 0x55)
COL_BAND = RGBColor(0xEA, 0xF1, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render_equation(latex_src: str, fontsize: int = 28, dpi: int = 300) -> str:
    """Render a LaTeX math string to a transparent PNG, return the path."""
    key = hashlib.md5(f"{latex_src}|{fontsize}|{dpi}".encode()).hexdigest()
    path = os.path.join(EQ_CACHE, f"{key}.png")
    if os.path.exists(path):
        return path
    fig = plt.figure(figsize=(0.01, 0.01))
    # Wrap in display-style for bigger rendering.
    fig.text(
        0.5, 0.5, f"${latex_src}$", fontsize=fontsize, ha="center", va="center"
    )
    plt.savefig(
        path, bbox_inches="tight", pad_inches=0.1, dpi=dpi, transparent=True
    )
    plt.close(fig)
    return path


def eq_dims(path, max_w_in=10.5, max_h_in=3.2):
    """Return (width, height) in Emu that fit within bounds preserving aspect."""
    from PIL import Image

    with Image.open(path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    w_in = min(max_w_in, max_h_in * aspect)
    h_in = w_in / aspect
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in * aspect
    return Inches(w_in), Inches(h_in)


def add_rect(slide, x, y, w, h, fill_rgb, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=16,
    bold=False,
    italic=False,
    color=COL_BODY,
    font="Calibri",
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(
    slide, items, x, y, w, h, *, size=14, color=COL_BODY, bullet="• "
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet}{item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_equation(slide, latex_src, cx, cy, *, fontsize=28, max_w=10.5,
                 max_h=3.0):
    path = render_equation(latex_src, fontsize=fontsize)
    w, h = eq_dims(path, max_w_in=max_w, max_h_in=max_h)
    x = cx - w // 2
    y = cy - h // 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)


def page_frame(slide, title, section=None):
    # top band
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.65), COL_TITLE)
    add_text(
        slide, title, Inches(0.45), Inches(0.08), Inches(11.5), Inches(0.55),
        size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        font="Calibri", anchor=MSO_ANCHOR.MIDDLE,
    )
    if section:
        add_text(
            slide, section, Inches(10.5), Inches(0.08), Inches(2.7), Inches(0.55),
            size=11, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC),
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )
    # bottom bar
    add_rect(slide, 0, Inches(7.35), SLIDE_W, Inches(0.15), COL_ACCENT)


def add_table(slide, data, x, y, w, h, *, header_fill=COL_TITLE,
              header_fg=RGBColor(0xFF, 0xFF, 0xFF), size=12,
              band=COL_BAND, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    shp = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            shp.columns[j].width = Emu(int(w * cw / total))
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = shp.cell(i, j)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = band
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            r = tf.paragraphs[0].add_run()
            r.text = str(val)
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.bold = i == 0
            r.font.color.rgb = header_fg if i == 0 else COL_BODY
    return shp


# ---------- Slide content ----------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(blank)


# ---- Slide 1: Title ----
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_TITLE)
add_rect(s, 0, Inches(7.0), SLIDE_W, Inches(0.5), COL_ACCENT)
add_text(
    s, "Toy-bubble models", Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2),
    size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), font="Calibri",
)
add_text(
    s, "e-NRTL → bubble thermodynamics in water electrolysis",
    Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.8),
    size=28, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC), font="Calibri",
)
add_text(
    s, "Mathematical framework for bubble_enrtl_toy.py",
    Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
    size=18, color=RGBColor(0xAE, 0xC6, 0xCF), font="Calibri",
)
add_text(
    s, "Nomenclature follows the Grok draft and standard electrolyte-thermodynamics literature.",
    Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
    size=14, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC), font="Calibri",
)


# ---- Slide 2: Pipeline overview ----
s = new_slide()
page_frame(s, "Four-layer pipeline", section="Overview")
layers = [
    ("A", "Symmetric-reference e-NRTL",
     "G^ex(T, P, {N_i}),  μ_i,  a_w,  γ_±"),
    ("B", "Butler surface-tension equation",
     "γ(m) = σ(m)"),
    ("C", "Axisymmetric Young–Laplace shape",
     "bubble profile,  V_d,  r_cl"),
    ("D", "VLE & nucleation in H₂ bubble",
     "y_H2,  P_bub,  r*"),
]
y = Inches(1.1)
row_h = Inches(1.3)
for i, (tag, name, out) in enumerate(layers):
    row_y = y + Inches(i * 1.4)
    add_rect(s, Inches(0.55), row_y, Inches(1.1), Inches(1.1), COL_TITLE)
    add_text(
        s, tag, Inches(0.55), row_y, Inches(1.1), Inches(1.1),
        size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, name, Inches(1.85), row_y + Inches(0.05), Inches(5.0), Inches(0.6),
        size=20, bold=True, color=COL_DARK, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, "→  " + out, Inches(6.95), row_y + Inches(0.05),
        Inches(6.0), Inches(1.0),
        size=16, italic=True, color=COL_ACCENT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ---- Slide 3: Starting definition ----
s = new_slide()
page_frame(s, "Thermodynamic definition of surface tension",
           section="Starting point  ·  Grok p. 3")
add_equation(
    s,
    r"\gamma^{\,j}\;=\;\left(\frac{\partial G^{\,j}}{\partial A^{\,j}}\right)_{T^{\,j},\,P^{\,j},\,N_i^{\,j}}",
    Inches(6.67), Inches(3.4), fontsize=44, max_h=2.4,
)
add_text(
    s,
    "Goal of Layer B: realise this derivative without a Langmuir fit, by placing the same e-NRTL model "
    "on both sides of the liquid–gas interface (Butler equation).",
    Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.5),
    size=17, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER,
)


# ---- Slide 4-5: Nomenclature (split in two) ----
nomen = [
    ("Symbol", "Meaning", "Units"),
    ("T, P", "temperature, pressure", "K, Pa"),
    ("Nᵢ, nᵢ", "amount (extensive, molar) of species i", "mol"),
    ("xᵢ", "mole fraction of species i", "–"),
    ("m", "salt molality (mol / kg water)", "mol kg⁻¹"),
    ("zᵢ", "signed charge of ion i", "–"),
    ("ν₊, ν₋, ν", "stoichiometric coefficients", "–"),
    ("I_x", "ionic strength, mole-fraction basis", "–"),
    ("I_m", "ionic strength, molality basis", "mol kg⁻¹"),
    ("Aφ", "Debye–Hückel parameter (water, 298 K) ≈ 0.392", "–"),
    ("ρ", "PDH closest-approach parameter = 14.9", "–"),
    ("α_ij", "NRTL non-randomness factor (fixed = 0.2)", "–"),
    ("τ_ij", "NRTL binary interaction parameter", "–"),
    ("G_ij = exp(–ατ)", "NRTL local-composition factor", "–"),
]
nomen2 = [
    ("Symbol", "Meaning", "Units"),
    ("G^ex", "total excess Gibbs energy", "J"),
    ("μᵢ", "chemical potential of species i", "J mol⁻¹"),
    ("aᵢ = xᵢ γᵢ", "activity of species i", "–"),
    ("γ_±", "mean ionic activity coefficient (molality basis)", "–"),
    ("a_w", "water activity", "–"),
    ("φ", "osmotic coefficient", "–"),
    ("γ = σ", "surface tension", "N m⁻¹"),
    ("Aᵢ", "partial molar surface area of species i", "m² mol⁻¹"),
    ("σᵢ⁰", "surface tension of (hypothetical) pure component i", "N m⁻¹"),
    ("k_H", "Henry's constant (solubility form), c = k_H p", "mol kg⁻¹ Pa⁻¹"),
    ("k_s", "Sechenov (salting-out) coefficient", "L mol⁻¹"),
    ("R_bub", "bubble radius", "m"),
    ("θ_g", "contact angle through the gas phase", "rad"),
    ("V_d, D_d", "detachment volume, diameter", "m³, m"),
    ("r*", "critical nucleation radius", "m"),
    ("S, S₀", "supersaturation (effective / reference)", "–"),
]
s = new_slide()
page_frame(s, "Nomenclature (part 1)", section="§0")
add_table(
    s, nomen, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.2),
    size=13, col_widths=[2, 6, 2.2],
)

s = new_slide()
page_frame(s, "Nomenclature (part 2)", section="§0")
add_table(
    s, nomen2, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.2),
    size=12, col_widths=[2, 6, 2.2],
)
add_text(
    s, "Subscripts:  c = cation,  a = anion,  w = water/solvent,  s = lumped salt,  "
       "B = bulk,  S = surface phase,  0 = reference / pure-water value.",
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
    size=11, italic=True, color=COL_MUTED,
)


# ---- Slide: §1 Layer A header ----
s = new_slide()
page_frame(s, "Layer A — Symmetric-reference e-NRTL", section="§1")
add_rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.4), COL_BAND)
add_bullets(s, [
    "Total molar Gibbs energy Gʲ = ideal + PDH + NRTL",
    "Long-range: Pitzer–Debye–Hückel (symmetric, mole-fraction)",
    "Short-range: local-composition NRTL (lumped-salt form)",
    "Chemical potentials from ∂G^ex/∂nᵢ (symmetric finite difference)",
    "Water activity a_w and osmotic coefficient φ",
    "Mean ionic activity coefficient γ_±(m) via Gibbs–Duhem",
], Inches(1.1), Inches(1.5), Inches(11.0), Inches(5.0), size=20,
   color=COL_DARK)


# ---- §1.1 Total molar Gibbs energy ----
s = new_slide()
page_frame(s, "§1.1  Total molar Gibbs energy of liquid phase j", section="Layer A")
add_text(
    s, "Modern symmetric convention (Chen / Song & Chen, Aspen-style):",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.5), size=16, color=COL_MUTED,
)
add_equation(
    s,
    r"G^{\,j}=\sum_i x_i\,G_i^{0}(T,P)+RT\sum_i x_i\,\ln(x_i\gamma_i)+G^{\text{ex},\,\text{PDH}}+G^{\text{ex},\,\text{NRTL}}",
    Inches(6.67), Inches(3.3), fontsize=30, max_h=2.0,
)
add_text(
    s,
    "Excess part decomposed into a long-range Pitzer–Debye–Hückel term "
    "and a short-range local-composition (NRTL) term.",
    Inches(0.6), Inches(5.8), Inches(12.0), Inches(1.0),
    size=15, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER,
)


# ---- §1.2 PDH ----
s = new_slide()
page_frame(s, "§1.2  Long-range Pitzer–Debye–Hückel term", section="Layer A")
add_text(
    s,
    "Symmetric reference on mole-fraction basis (Chen & Evans 1986; Song & Chen 2009):",
    Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.5),
    size=14, color=COL_MUTED,
)
add_equation(
    s,
    r"\frac{G^{\text{ex},\,\text{PDH}}}{n_t RT}=-\left(\frac{1000}{M_w}\right)^{1/2}\frac{4A_\phi I_x}{\rho}\ln\!\bigl(1+\rho\sqrt{I_x}\bigr)",
    Inches(6.67), Inches(2.3), fontsize=28, max_h=1.6,
)
add_equation(
    s,
    r"I_x=\tfrac12\sum_i z_i^{\,2}\,x_i,\qquad M_w=18.015\,\text{g mol}^{-1},\qquad A_\phi=0.392,\qquad \rho=14.9",
    Inches(6.67), Inches(4.0), fontsize=22, max_h=1.0,
)
add_bullets(s, [
    "Recovers DH limiting law  ln γ_± → –|z₊z₋| Aφ √I_m  as m → 0",
    "Responsible for the initial droop of γ_± (Fig. 1)",
    "Carries the charge structure → KOH vs H₂SO₄ diverge already at low m (Fig. 6)",
], Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.0), size=14, color=COL_BODY)


# ---- §1.3 NRTL ----
s = new_slide()
page_frame(s, "§1.3  Short-range local-composition NRTL term", section="Layer A")
add_text(
    s,
    "Single strong electrolyte + solvent with lumped-salt symmetric form; X_s = x_c + x_a:",
    Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.5),
    size=14, color=COL_MUTED,
)
add_equation(
    s,
    r"\frac{G^{\text{ex},\,\text{NRTL}}}{n_t RT}=x_w\,X_s\left[\frac{\tau_{sw}G_{sw}}{x_w+X_s G_{sw}}+\frac{\tau_{ws}G_{ws}}{X_s+x_w G_{ws}}\right]",
    Inches(6.67), Inches(2.7), fontsize=28, max_h=1.8,
)
add_bullets(s, [
    "G_{ij} = exp(–α τ_{ij}),  α = 0.2",
    "Two parameters per salt: τ_{sw} (salt→solvent) and τ_{ws} (solvent→salt)",
    "Like-ion repulsion enforced by lumping ion pair → τ_{cc} = τ_{aa} = 0 by construction",
], Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.0), size=15, color=COL_BODY)


# ---- §1.3 Caveat ----
s = new_slide()
page_frame(s, "§1.3  Toy-level caveat on the lumped-salt form", section="Layer A")
add_rect(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(5.8), COL_BAND)
add_text(
    s, "Why the lumped form suffices here",
    Inches(0.9), Inches(1.1), Inches(11.5), Inches(0.5),
    size=20, bold=True, color=COL_TITLE,
)
add_bullets(s, [
    "Full Chen symmetric e-NRTL keeps cation & anion as separate interacting species, with "
    "charge-weighted effective mole fractions X_i = |z_i| x_i for ions.",
    "That gives slightly different short-range non-ideality for asymmetric salts (e.g. H₂SO₄).",
    "We adopt the lumped form because:",
    "   (i) the PDH part already carries the dominant charge-asymmetry at low–moderate m.",
    "   (ii) ion-pair parameters for single-solvent binaries are effectively indistinguishable in "
    "their observable consequences (γ_±, a_w, φ) — they would just recalibrate.",
], Inches(1.0), Inches(1.7), Inches(11.5), Inches(5.0), size=15, color=COL_DARK)


# ---- §1.4 Chemical potentials ----
s = new_slide()
page_frame(s, "§1.4  Chemical potentials and activity coefficients", section="Layer A")
add_equation(
    s,
    r"\frac{\mu_i^{\text{ex}}}{RT}=\left.\frac{\partial}{\partial n_i}\!\left(\frac{G^{\text{ex}}}{RT}\right)\right|_{T,P,\,n_{k\neq i}}",
    Inches(6.67), Inches(2.2), fontsize=28, max_h=1.8,
)
add_equation(
    s,
    r"a_i=x_i\gamma_i=x_i\,\exp\!\bigl(\mu_i^{\text{ex}}/RT\bigr)",
    Inches(6.67), Inches(4.3), fontsize=30, max_h=1.2,
)
add_text(
    s,
    "In the code these partials are computed by symmetric finite difference on the extensive G^ex "
    "(chemical_potentials in bubble_enrtl_toy.py). Sidesteps hand-derivation errors; "
    "stable well below 1 mol/kg.",
    Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.2),
    size=14, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER,
)


# ---- §1.5 Water activity & osmotic coeff ----
s = new_slide()
page_frame(s, "§1.5  Water activity and osmotic coefficient", section="Layer A")
add_text(
    s,
    "At molality m with  n_w = 1000 / M_w  mol water per kg:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.5), size=15, color=COL_MUTED,
)
add_equation(
    s,
    r"a_w(m)=x_w\,\exp\!\bigl(\mu_w^{\text{ex}}/RT\bigr),\qquad \phi(m)=-\frac{\ln a_w(m)}{M_w^{[\text{kg}]}\,\nu\,m}",
    Inches(6.67), Inches(3.4), fontsize=26, max_h=1.6,
)


# ---- §1.6 Gibbs-Duhem ----
s = new_slide()
page_frame(s, "§1.6  Mean ionic activity coefficient via Gibbs–Duhem", section="Layer A")
add_text(
    s,
    "To avoid the unsymmetric-to-molality reference-state transformation, integrate the solvent activity:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.5), size=14, color=COL_MUTED,
)
add_equation(
    s,
    r"\ln\gamma_\pm(m)\;=\;\bigl(\phi(m)-1\bigr)+\int_0^{m}\frac{\phi(m')-1}{m'}\,\mathrm d m'",
    Inches(6.67), Inches(3.2), fontsize=28, max_h=1.8,
)
add_bullets(s, [
    "Gives γ_± on the standard molality basis (pure water at infinite dilution = reference)",
    "Integral evaluated with scipy.integrate.cumulative_trapezoid on a 40-point molality grid",
], Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6), size=15, color=COL_BODY)


# ---- Layer B header ----
s = new_slide()
page_frame(s, "Layer B — Butler surface-tension equation", section="§2")
add_equation(
    s,
    r"\gamma^{\,j}=\left(\frac{\partial G^{\,j}}{\partial A^{\,j}}\right)_{T,P,\{N_i\}}",
    Inches(6.67), Inches(2.2), fontsize=34, max_h=1.8,
)
add_text(
    s,
    "For a flat liquid–gas interface in LTE with the bulk, the chemical potential of each species "
    "equals in the surface phase (S) and bulk (B). One Butler equation per species:",
    Inches(0.6), Inches(3.6), Inches(12.0), Inches(1.0), size=15, color=COL_MUTED,
    align=PP_ALIGN.CENTER,
)
add_equation(
    s,
    r"\boxed{\;\sigma=\sigma_i^{0}+\dfrac{RT}{A_i}\,\ln\!\dfrac{a_i^{S}}{a_i^{B}}\qquad(\text{each }i)\;}",
    Inches(6.67), Inches(5.7), fontsize=28, max_h=1.5,
)


# ---- §2.2 Butler system solve ----
s = new_slide()
page_frame(s, "§2.2  Butler equation — system solve", section="Layer B")
add_bullets(s, [
    "For binary water + lumped-salt: a 2×2 system in (σ, x_s^S)",
    "σ eliminated by equating the water and salt equations → one nonlinear equation in x_s^S",
    "Solved with scipy.optimize.brentq in butler_surface_tension",
    "Reference: Butler, Proc. R. Soc. A 135, 348 (1932)",
], Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.6), size=18, color=COL_DARK)
add_text(
    s,
    "§2.3  Activity coefficients at the surface",
    Inches(0.6), Inches(4.2), Inches(12.0), Inches(0.5),
    size=20, bold=True, color=COL_TITLE,
)
add_bullets(s, [
    "Li & Lu family (Langmuir 17, 3532, 2001; Hu & Lee FPE 158, 1043, 1999) couple surface phase to its OWN e-NRTL at surface composition.",
    "TOY approximation: γ_i^S(x^S) = γ_i^B(x = x^S) — the same e-NRTL evaluated at the surface composition.",
    "Preserves the Butler structure and gets the trend right; full surface-phase recalibration is a later refinement.",
], Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.4), size=13, color=COL_BODY)


# ---- §2.4 Partial molar surface area ----
s = new_slide()
page_frame(s, "§2.4  Partial molar surface area", section="Layer B")
add_text(
    s,
    "Stefan–Guggenheim / Sprow–Prausnitz form:",
    Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.5), size=16, color=COL_MUTED,
)
add_equation(
    s,
    r"A_i=f\,N_A^{\,1/3}\,V_i^{\,2/3},\qquad f\approx 1.091",
    Inches(6.67), Inches(3.0), fontsize=32, max_h=1.6,
)
add_text(
    s,
    "V_i  =  partial molar volume.   Numerical values used in the toy are tabulated in §5.2.",
    Inches(0.6), Inches(5.2), Inches(12.0), Inches(1.0),
    size=16, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER,
)


# ---- Layer C header ----
s = new_slide()
page_frame(s, "Layer C — Axisymmetric Young–Laplace shape & detachment",
           section="§3")
add_rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.4), COL_BAND)
add_bullets(s, [
    "Young–Laplace ODE in arc-length parameterisation (from the apex)",
    "Contact-angle boundary condition through the gas phase",
    "Fritz-type quasi-static detachment force balance",
    "Solve for detachment volume V_d(σ, θ_g) and diameter D_d",
], Inches(1.1), Inches(1.6), Inches(11.0), Inches(5.0), size=22, color=COL_DARK)


# ---- §3.1 Young-Laplace ODE ----
s = new_slide()
page_frame(s, "§3.1  Young–Laplace ODE (sessile bubble)", section="Layer C")
add_text(
    s,
    "Bubble on horizontal solid at z = 0, apex at z = H, buoyancy lifting upward. "
    "Arc-length s parameterisation from the apex:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.7), size=13, color=COL_MUTED,
)
add_equation(
    s,
    r"\frac{\mathrm d r}{\mathrm d s}=\cos\varphi,\qquad \frac{\mathrm d z_d}{\mathrm d s}=\sin\varphi,\qquad \frac{\mathrm d\varphi}{\mathrm d s}=\frac{2}{b}-\beta z_d-\frac{\sin\varphi}{r}",
    Inches(6.67), Inches(3.0), fontsize=24, max_h=1.6,
)
add_bullets(s, [
    "z_d = H – z  (depth below apex);  b = apex radius of curvature (shooting parameter)",
    "β = (ρ_l – ρ_g) g / σ",
    "IC:  r(0)=0,  z_d(0)=0,  φ(0)=0",
    "Near-apex seeding with spherical expansion  r ≈ b sin(s/b),  z_d ≈ b(1–cos(s/b))  "
    "to avoid the r=0 singularity",
], Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), size=13, color=COL_BODY)


# ---- §3.2 Contact angle BC ----
s = new_slide()
page_frame(s, "§3.2  Contact-angle boundary condition", section="Layer C")
add_text(
    s,
    "Contact angle θ_g measured THROUGH the gas phase at the three-phase line. "
    "The contact line is reached when the tangent angle satisfies:",
    Inches(0.6), Inches(1.1), Inches(12.0), Inches(1.0), size=14, color=COL_MUTED,
)
add_equation(
    s,
    r"\varphi_{\rm cl}=\pi-\theta_g",
    Inches(6.67), Inches(3.3), fontsize=36, max_h=1.5,
)
add_bullets(s, [
    "Integration terminates on this event (event_phi in young_laplace_shape)",
    "Secondary event catches pinch-off  (r → 0 again)",
], Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8), size=15, color=COL_BODY)


# ---- §3.3 Fritz ----
s = new_slide()
page_frame(s, "§3.3  Fritz-type detachment force balance", section="Layer C")
add_text(
    s,
    "Quasi-static detachment: buoyancy lifting the bubble equals the vertical component of surface "
    "tension at the pinning line",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.9),
    size=14, color=COL_MUTED,
)
add_equation(
    s,
    r"(\rho_l-\rho_g)\,g\,V_{\rm bub}=2\pi r_{\rm cl}\,\sigma\,\sin\theta_g",
    Inches(6.67), Inches(3.0), fontsize=30, max_h=1.6,
)
add_text(
    s, "(Fritz, Phys. Z. 36, 379, 1935)",
    Inches(0.6), Inches(4.2), Inches(12.0), Inches(0.4), size=12, italic=True,
    color=COL_MUTED, align=PP_ALIGN.CENTER,
)
add_bullets(s, [
    "Rooted on the shooting parameter b;  detachment_volume scans b over a multiple of the capillary "
    "length  ℓ_c = √(σ / Δρ g)",
    "Brackets the sign change and refines with brentq",
    "Output:  V_d(σ, θ_g)  and  D_d = (6 V_d / π)¹ᐟ³",
], Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.2), size=13, color=COL_BODY)


# ---- Layer D header ----
s = new_slide()
page_frame(s, "Layer D — VLE and nucleation", section="§4")
add_rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.4), COL_BAND)
add_bullets(s, [
    "Effective Henry's constant with salting-out (Sechenov OR a_w form)",
    "Bubble composition: Raoult for water + Laplace for P",
    "Critical nucleation radius r*(m) from Young–Laplace balance",
], Inches(1.1), Inches(1.6), Inches(11.0), Inches(5.0), size=22, color=COL_DARK)


# ---- §4.1 Henry ----
s = new_slide()
page_frame(s, "§4.1  Effective Henry's constant (salting-out)", section="Layer D")
add_text(
    s, "Two standard options, both produce the same trend in the toy:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.5), size=15, color=COL_MUTED,
)
add_equation(
    s,
    r"\text{Sechenov:}\quad k_H^{\text{eff}}(m)=k_H^{0}\cdot 10^{-k_s\,c_{\rm salt}}",
    Inches(6.67), Inches(2.4), fontsize=26, max_h=1.2,
)
add_equation(
    s,
    r"\text{Water-activity (KK-like):}\quad k_H^{\text{eff}}(m)\approx k_H^{0}\big/a_w(m)",
    Inches(6.67), Inches(4.0), fontsize=26, max_h=1.2,
)
add_bullets(s, [
    "Sechenov uses tabulated k_s values",
    "Water-activity form uses a_w from the same e-NRTL — no extra parameter beyond §1",
    "Either route shifts solubility of a neutral gas (H₂) through the non-ideal water",
], Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.8), size=13, color=COL_BODY)


# ---- §4.2 Bubble composition ----
s = new_slide()
page_frame(s, "§4.2  Bubble composition  (H₂ + H₂O vapour)", section="Layer D")
add_text(
    s, "Raoult with activity for water + Laplace for total pressure:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.5),
    size=16, color=COL_MUTED,
)
add_equation(
    s,
    r"p_{\rm H_2O}=a_w(m)\,p_{\rm sat}(T),\qquad P_{\rm bub}=P_{\rm atm}+\frac{2\sigma(m)}{R_{\rm bub}},\qquad y_{\rm H_2}=\frac{P_{\rm bub}-p_{\rm H_2O}}{P_{\rm bub}}",
    Inches(6.67), Inches(3.6), fontsize=24, max_h=2.0,
)


# ---- §4.3 r* ----
s = new_slide()
page_frame(s, "§4.3  Critical nucleation radius", section="Layer D")
add_text(
    s,
    "At fixed dissolved-H₂ corresponding to reference supersaturation S₀ in pure water, "
    "the effective local supersaturation in salty electrolyte is:",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.8), size=13, color=COL_MUTED,
)
add_equation(
    s,
    r"S_{\rm eff}(m)=S_0\,\frac{k_H^{0}}{k_H^{\text{eff}}(m)}",
    Inches(6.67), Inches(2.6), fontsize=28, max_h=1.4,
)
add_text(
    s,
    "Young–Laplace critical radius:",
    Inches(0.6), Inches(3.7), Inches(12.0), Inches(0.5), size=14, color=COL_MUTED,
    align=PP_ALIGN.CENTER,
)
add_equation(
    s,
    r"r^{*}(m)=\frac{2\,\sigma(m)}{(S_{\rm eff}-1)\,P_{\rm atm}}",
    Inches(6.67), Inches(4.9), fontsize=30, max_h=1.6,
)
add_text(
    s,
    "The compound behaviour of σ(m) ↑ and S_eff(m) ↑ pulls r* down sharply as m rises "
    "(Fig. 5 right panel).",
    Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.7),
    size=13, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER,
)


# ---- §5 Parameters ----
s = new_slide()
page_frame(s, "§5.1  e-NRTL parameters", section="Parameters")
enrtl_tbl = [
    ("Salt", "z_c", "z_a", "ν_c", "ν_a", "τ_sw", "τ_ws", "α"),
    ("KOH", "+1", "−1", "1", "1", "11.5", "−4.5", "0.2"),
    ("H₂SO₄", "+1", "−2", "2", "1", "13.0", "−5.2", "0.2"),
]
add_table(s, enrtl_tbl, Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.7),
          size=18)
add_text(
    s,
    "Values are representative — tuned only to give qualitatively correct γ_±(m) and a_w(m) shapes:\n"
    "   •  DH droop at small m\n"
    "   •  NRTL rebound at large m\n\n"
    "Published Aspen-style parameters for KOH–H₂O (Thomsen & Rasmussen, Chem. Eng. Sci. 54, 1787, 1999) "
    "should replace these once quantitative fidelity is wanted.",
    Inches(0.7), Inches(3.5), Inches(12.0), Inches(3.3),
    size=15, color=COL_MUTED,
)


s = new_slide()
page_frame(s, "§5.2  Butler parameters", section="Parameters")
butler_tbl = [
    ("Species", "σᵢ⁰ (N/m)", "Aᵢ (m² mol⁻¹)", "Vᵢ (m³ mol⁻¹)"),
    ("H₂O", "0.07197", "6.30×10⁴", "1.80×10⁻⁵"),
    ("KOH", "0.170", "8.40×10⁴", "2.70×10⁻⁵"),
    ("H₂SO₄", "0.155", "1.31×10⁵", "5.30×10⁻⁵"),
]
add_table(s, butler_tbl, Inches(1.0), Inches(1.4), Inches(11.3), Inches(2.6),
          size=18, col_widths=[2, 2, 3, 3])


s = new_slide()
page_frame(s, "§5.3  Henry, salting-out, nucleation  ·  §5.4  General", section="Parameters")
hs_tbl = [
    ("Parameter", "Value", "Note"),
    ("k_H⁰ (H₂, 25°C)", "7.8×10⁻⁹ mol kg⁻¹ Pa⁻¹",
     "Sander, Atmos. Chem. Phys. 15, 4399 (2015)"),
    ("k_s (KOH)", "0.134 L mol⁻¹", "≈ Sechenov coefficient for H₂"),
    ("k_s (H₂SO₄)", "0.099 L mol⁻¹", "≈ Sechenov coefficient for H₂"),
    ("θ_g", "30°", "hydrophobic cavity, Raman-setup-like"),
    ("S₀", "2.0", "reference supersaturation in pure water"),
]
add_table(s, hs_tbl, Inches(0.4), Inches(0.95), Inches(7.8), Inches(3.6),
          size=12, col_widths=[2.4, 3.0, 4.0])

add_text(s, "§5.4  General conditions",
         Inches(8.4), Inches(0.95), Inches(4.6), Inches(0.4),
         size=16, bold=True, color=COL_TITLE)
gen_rows = [
    ("Parameter", "Value"),
    ("T", "298.15 K"),
    ("P_atm", "1.01325×10⁵ Pa"),
    ("ρ_l (water, 25°C)", "997.05 kg m⁻³"),
    ("ρ_g (H₂, 25°C, 1 bar)", "0.0899 kg m⁻³"),
    ("p_sat(T), 25°C", "3169 Pa"),
]
add_table(s, gen_rows, Inches(8.4), Inches(1.4), Inches(4.6), Inches(3.2),
          size=12, col_widths=[2.4, 2.2])

add_text(
    s,
    "These constants define the reference state used throughout the toy. "
    "T and P_atm fix the e-NRTL and Butler evaluations; ρ_l, ρ_g feed the "
    "Young–Laplace shape (β = (ρ_l − ρ_g) g / σ) and the Fritz detachment "
    "balance; p_sat(T) sets the water partial pressure in the bubble (§4.2).",
    Inches(0.4), Inches(5.0), Inches(12.6), Inches(2.0),
    size=14, italic=True, color=COL_MUTED,
)


# ---- §6 Salt vs ions vs field (1) ----
s = new_slide()
page_frame(s, "§6  Salt vs. ions vs. electric field — what the toy captures", section="Discussion")
add_text(
    s,
    "Common confusion: “you're varying salt molality, not ion concentration or electric field — "
    "aren't those the bigger effects?”  For a single, fully dissociated strong electrolyte these three "
    "are NOT independent.",
    Inches(0.6), Inches(1.0), Inches(12.0), Inches(1.3), size=14,
    italic=True, color=COL_MUTED,
)
add_bullets(s, [
    "Ion concentration is set by salt molality via stoichiometry: c₊ = ν₊m,  c₋ = ν₋m. "
    "KOH at 3 mol/kg means 3 mol/kg K⁺ AND 3 mol/kg OH⁻. Varying m IS varying ion concentration.",
    "Ionic strength: I_m = ½ Σ νᵢ zᵢ² m. For KOH, I_m = m; for H₂SO₄, I_m = 3m. "
    "PDH (§1.2) captures this exactly — that is why KOH and H₂SO₄ diverge in Fig. 1.",
    "Ion-identity effects beyond ionic strength come from the e-NRTL short-range τ parameters. "
    "Plotting σ vs I_m does NOT collapse the two salts onto one curve (Fig. 6 right panel) — "
    "specific-ion behaviour is real and captured by e-NRTL.",
], Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.5), size=14, color=COL_BODY)


# ---- §6 Salt vs ions vs field (2) ----
s = new_slide()
page_frame(s, "§6  Electric-field / double-layer effects are NOT in this toy", section="Discussion")
add_bullets(s, [
    "They act at the ELECTRODE–electrolyte interface (polarised Pt), through:",
    "   •  Gouy–Chapman–Stern diffuse-layer structure",
    "   •  Lippmann equation,  ∂σ/∂E = –q_s,  linking electrode surface charge to electrocapillary depression",
    "   •  Frumkin correction to activation overpotential",
    "The bubble–electrolyte interface here is NOT polarised in the same way.",
    "It carries a weak ζ-potential (typically −30 to −50 mV for H₂ bubbles in water) from "
    "preferential ion adsorption — implicit in the Butler surface-phase composition.",
], Inches(0.7), Inches(1.1), Inches(11.8), Inches(4.0), size=14, color=COL_BODY)
add_text(s,
         "A full treatment of electric-field effects on σ would add a Lippmann term:",
         Inches(0.7), Inches(5.1), Inches(12.0), Inches(0.5),
         size=14, color=COL_MUTED)
add_equation(
    s,
    r"\sigma(E)=\sigma_0-\tfrac12 C_{\rm dl}(E-E_{\rm pzc})^2",
    Inches(6.67), Inches(6.1), fontsize=24, max_h=1.0,
)


# ---- §7 Limitations ----
s = new_slide()
page_frame(s, "§7  Known limitations and extensions (1/2)", section="Limitations")
lim_tbl_1 = [
    ("#", "Simplification", "Path to improvement"),
    ("1", "“Lumped-salt” NRTL (not full Chen-symmetric)",
     "Full Chen 1986 / Aspen symmetric e-NRTL (two τ pairs per ion-solvent pair)"),
    ("2", "Bulk e-NRTL used for surface activities in Butler",
     "Couple a surface-phase e-NRTL as in Li & Lu 2001"),
    ("3", "τ parameters chosen by trend, not fitted to data",
     "Published Aspen parameters (Thomsen & Rasmussen 1999 KOH–H₂O; Clegg & Brimblecombe 1995 H₂SO₄–H₂O)"),
    ("4", "Fritz-style vertical-force detachment",
     "Pinch-off from the actual YL shape (watch r_neck → 0)"),
    ("5", "Quasi-static bubble (no inertia / drag / Marangoni)",
     "Couple σ(m) into dynamic integrators in ddgclib/dynamic_integrators/"),
]
add_table(s, lim_tbl_1, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.8),
          size=12, col_widths=[0.6, 5, 7.5])


s = new_slide()
page_frame(s, "§7  Known limitations and extensions (2/2)", section="Limitations")
lim_tbl_2 = [
    ("#", "Simplification", "Path to improvement"),
    ("6", "No double-layer / electrocapillarity",
     "Add Lippmann correction to σ at polarised surface"),
    ("7", "Sechenov form for salting-out",
     "Derive salting-out directly from e-NRTL activity coefficient of dissolved H₂ "
     "(needs τ parameters for H₂–water and H₂–ion)"),
    ("8", "Ideal-gas bubble (no H₂ fugacity correction)",
     "Peng–Robinson or similar EoS for φ_H₂(T,P)"),
    ("9", "Single solvent (water)",
     "Extend to mixed solvents (water + methanol etc. for CO₂ electrolysis)"),
    ("10", "Constant T, P",
     "Add T, P dependence of all parameters (heavy but mechanical)"),
]
add_table(s, lim_tbl_2, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.8),
          size=12, col_widths=[0.6, 5, 7.5])


# ---- §8 References (primary) ----
s = new_slide()
page_frame(s, "§8  Primary literature the toy comments on", section="References")
refs = [
    ("Sepahi et al. (2022) — the “Krug paper”",
     "The effect of buoyancy-driven convection on the growth and dissolution of bubbles on "
     "electrodes. Electrochim. Acta 403, 139616.  doi:10.1016/j.electacta.2021.139616",
     "linear β, ideal gas, constant k_H, no composition-dependent σ"),
    ("Raman et al. (2022) — the “Rivas paper”",
     "Potential response of single successive constant-current-driven electrolytic hydrogen "
     "bubbles spatially separated from the electrode. Electrochim. Acta 425, 140691.  "
     "doi:10.1016/j.electacta.2022.140691",
     "fixed σ = 72 mN/m; concentration overpotential inferred by subtraction"),
]
y = 1.1
for title, body, tag in refs:
    add_rect(s, Inches(0.5), Inches(y), Inches(0.15), Inches(2.7), COL_ACCENT)
    add_text(s, title, Inches(0.8), Inches(y), Inches(11.8), Inches(0.45),
             size=17, bold=True, color=COL_DARK)
    add_text(s, body, Inches(0.8), Inches(y + 0.45), Inches(11.8), Inches(1.3),
             size=13, color=COL_BODY)
    add_text(s, "→  " + tag, Inches(0.8), Inches(y + 1.7), Inches(11.8),
             Inches(0.5), size=13, italic=True, color=COL_ACCENT)
    y += 3.0


# ---- §8 References (foundations) ----
s = new_slide()
page_frame(s, "§8  Thermodynamic model foundations", section="References")
foundations = [
    ("Butler, J. A. V.", "Proc. R. Soc. A 135, 348 (1932)",
     "Butler equation for surface tension"),
    ("Chen, C.-C., Evans, L. B.", "AIChE J. 32, 444 (1986)",
     "symmetric e-NRTL (local composition, excess Gibbs energy)"),
    ("Song, Y., Chen, C.-C.", "Ind. Eng. Chem. Res. 48, 7788 (2009)",
     "modern symmetric form"),
    ("Li, Z.-B., Lu, B. C.-Y.", "Langmuir 17, 3532 (2001)",
     "Butler + e-NRTL coupling, aqueous electrolyte surface tensions"),
    ("Hu, Y.-F., Lee, H.", "Fluid Phase Equil. 158-160, 1043 (1999)",
     "surface-phase formulation"),
    ("Thomsen, K., Rasmussen, P.", "Chem. Eng. Sci. 54, 1787 (1999)",
     "published KOH–H₂O Aspen-style parameters"),
    ("Fritz, W.", "Phys. Z. 36, 379 (1935)",
     "detachment force balance"),
    ("Sechenov, I. M.", "Z. Phys. Chem. 4, 117 (1889)",
     "salting-out linear correlation"),
    ("Sander, R.", "Atmos. Chem. Phys. 15, 4399 (2015)",
     "compilation of Henry's law constants (H₂, water)"),
]
ref_rows = [("Author", "Reference", "Topic")] + list(foundations)
add_table(s, ref_rows, Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.9),
          size=12, col_widths=[2.3, 3.2, 6.5])


# ---- Save ----
prs.save(OUT_PPTX)
print(f"Wrote {OUT_PPTX}")
print(f"Slides: {len(prs.slides)}")
